import sys
import tkinter
from pptx import Presentation


file = open('Gospel_Lyrics.txt')
AllSong = file.read()
song = AllSong.split('/')
title = []
title_h = []
song_sp = []
i = 0
for di in song:
    song_sp = di.split(',')
    title.insert(i, song_sp[0])
    title_h.insert(i, song_sp[1])
    print(title[i])
    print(title_h[i])
    i = i + 1

root = tkinter.Tk()
root.title(u"Gospel Insert Power Point")
root.geometry("400x300")

# 予測変換
def conversion(event):
    li_prediction.delete(0, tkinter.END)
    title = tb_title.get()
    i = 0
    for t in title_h:
        if title in t:
            li_prediction.insert(i, title_h[i])
        i = i + 1
    print(title)

# ひらがなの曲名から漢字の曲名に変換(convert)
def namecon(name):
    i = 0
    for na in title_h:
        if na == name:
            return title[i]
        i = i + 1

def numcon(name):
    i = 0
    for na in title:
        if na == name:
            return i
        i = i + 1

# 予測変換のリストボックスから曲順リストボックスに挿入
def insert(event):
    for i in li_prediction.curselection():
        name = li_prediction.get(i)
        li_order.insert("end", namecon(name))

# 選択した要素を一番上に移動
def totop(event):
    for i in li_order.curselection():
        name = li_order.get(i)
        li_order.delete(i)
        li_order.insert(0, name)

# パワーポイントを出力
def Create(event):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    flg = []
    for i in range(li_order.size()):
        t = li_order.get(i)
        slide = prs.slides.add_slide(layout)
        s_title = slide.placeholders[0]
        s_title.text = t

        # 何故かrstripが使えない
        
        # s_title.text = str(t).rstrip('\n')
        songNum = numcon(t)
        oneSong = song[songNum]
        flg = oneSong.split(',')
        for j in range(2, len(flg)):
            lyrics = slide.placeholders[1]
            lyrics.text = flg[j]
            # 曲と曲の間に白紙のスライドを挟みたいなら、下のif文を消す
            if j < len(flg) - 1:
                slide = prs.slides.add_slide(layout)
        print(flg[2])

    prs.save('test.pptx')

    # 曲名のみをスライドの最初に持ってくるか、曲名と歌詞を持ってくるか
    # for title in li_order:
    #     if isTitle:
    #         sld_layout = 0
    #     else:
    #         sld_layout = 1
    #     slide_layout = prs.slide_layout[sld_layout]
    
    pass

pw_main = tkinter.PanedWindow(orient='horizontal')
pw_main.pack(expand=True, fill = tkinter.BOTH, side="left")

pw_left = tkinter.PanedWindow(pw_main, orient='vertical')
pw_main.add(pw_left)
pw_right = tkinter.PanedWindow(pw_main, orient='vertical')
pw_main.add(pw_right)

lb_title = tkinter.Label(pw_left, text=u'タイトルを入力してください')
tb_title = tkinter.Entry(pw_left, width=25)
lb_title.pack(padx=20, pady=15)
tb_title.pack()

bt_prediction = tkinter.Button(pw_left, text=u'予測変換')
bt_prediction.bind("<Button-1>", conversion)
li_prediction = tkinter.Listbox(pw_left, height=8, width=25)
bt_prediction.pack(pady=10)
li_prediction.pack()

bt_insert = tkinter.Button(pw_left, text=u"Insert", width=10)
bt_insert.bind("<Button-1>", insert)
bt_insert.pack(pady=10)

li_order = tkinter.Listbox(pw_right, height=12, width=30)
li_order.grid(row=0, column=0, columnspan=2, padx=15, pady=20)

bt_up = tkinter.Button(pw_right, text=u"To the Top", width=10)
bt_up.bind("<Button-1>", totop)
bt_up.grid(row=1, column=0)

bt_create = tkinter.Button(pw_right, text=u"Create", width=10)
bt_create.bind("<Button-1>", Create)
bt_create.grid(row=1, column=1)

root.mainloop()
